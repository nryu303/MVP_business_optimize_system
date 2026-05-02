# -*- coding: utf-8 -*-
"""
AlphaScope案件 提案PPT生成
佐藤大将様向け 投資情報分析AI MVP提案
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT_PATH = r"c:\MVP\milestone1\AlphaScope_提案資料.pptx"

# ===== カラーパレット =====
NAVY = RGBColor(0x1F, 0x4E, 0x78)
BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT = RGBColor(0xDD, 0xEB, 0xF7)
GRAY = RGBColor(0x59, 0x59, 0x59)
DARK = RGBColor(0x1F, 0x1F, 0x1F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xED, 0x7D, 0x31)
GREEN = RGBColor(0x70, 0xAD, 0x47)

FONT_JP = "Yu Gothic"

# ===== プレゼン初期化 (16:9) =====
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_JP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(45000)
    tf.margin_right = Emu(45000)
    tf.margin_top = Emu(20000)
    tf.margin_bottom = Emu(20000)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(text, str):
        text_lines = text.split("\n")
    else:
        text_lines = text
    for i, line in enumerate(text_lines):
        if i == 0:
            run = p.add_run()
        else:
            np = tf.add_paragraph()
            np.alignment = align
            run = np.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_header(slide, title, subtitle=None):
    """各スライド上部のヘッダ"""
    add_rect(slide, 0, 0, SW, Inches(0.9), NAVY)
    add_rect(slide, 0, Inches(0.9), SW, Inches(0.06), ACCENT)
    add_text(slide, Inches(0.4), Inches(0.15), SW - Inches(0.8), Inches(0.55),
             title, size=22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.4), Inches(0.55), SW - Inches(0.8), Inches(0.35),
                 subtitle, size=11, color=LIGHT)


def add_footer(slide, page_no, total):
    add_text(slide, Inches(0.3), SH - Inches(0.35), Inches(8), Inches(0.3),
             "AlphaScope - 投資情報分析AI MVP 提案資料 / 提案者: daiske",
             size=9, color=GRAY)
    add_text(slide, SW - Inches(1.2), SH - Inches(0.35), Inches(0.9), Inches(0.3),
             f"{page_no} / {total}", size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def add_table(slide, x, y, w, h, headers, rows,
              col_ratios=None, row_height=None,
              header_fill=BLUE, header_font_color=WHITE,
              body_font_size=11, header_font_size=12):
    """シンプルなテーブル描画 (shapesベース)"""
    n_cols = len(headers)
    n_rows = len(rows) + 1
    if col_ratios is None:
        col_ratios = [1] * n_cols
    total_ratio = sum(col_ratios)
    col_widths = [int(w * r / total_ratio) for r in col_ratios]
    if row_height is None:
        row_height = h // n_rows
    cur_x = x
    # ヘッダ
    for ci, head in enumerate(headers):
        add_rect(slide, cur_x, y, col_widths[ci], row_height, header_fill, line=WHITE)
        add_text(slide, cur_x, y, col_widths[ci], row_height,
                 head, size=header_font_size, bold=True, color=header_font_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cur_x += col_widths[ci]
    # ボディ
    for ri, row in enumerate(rows):
        cur_x = x
        ry = y + row_height * (ri + 1)
        fill = WHITE if ri % 2 == 0 else LIGHT
        for ci, val in enumerate(row):
            add_rect(slide, cur_x, ry, col_widths[ci], row_height, fill,
                     line=RGBColor(0xBF, 0xBF, 0xBF))
            align = PP_ALIGN.CENTER if (ci > 0 and isinstance(val, str)
                                        and len(val) <= 4) else PP_ALIGN.LEFT
            add_text(slide, cur_x, ry, col_widths[ci], row_height, str(val),
                     size=body_font_size, color=DARK, align=align,
                     anchor=MSO_ANCHOR.MIDDLE)
            cur_x += col_widths[ci]


# ============================================================
# Slide 1: 表紙
# ============================================================
slide = add_slide()
add_rect(slide, 0, 0, SW, SH, NAVY)
# アクセントバー
add_rect(slide, 0, Inches(2.6), SW, Inches(0.08), ACCENT)
add_rect(slide, 0, Inches(5.2), SW, Inches(0.04), BLUE)

add_text(slide, Inches(0.8), Inches(2.85), SW - Inches(1.6), Inches(1.2),
         "AlphaScope", size=64, bold=True, color=WHITE)
add_text(slide, Inches(0.8), Inches(3.85), SW - Inches(1.6), Inches(0.8),
         "AI投資情報分析SaaS - MVP開発 提案資料",
         size=24, color=LIGHT)
add_text(slide, Inches(0.8), Inches(5.5), SW - Inches(1.6), Inches(0.5),
         "佐藤大将 様",
         size=18, color=WHITE)
add_text(slide, Inches(0.8), Inches(5.95), SW - Inches(1.6), Inches(0.4),
         "提案者: daiske    /    2026-04-30",
         size=14, color=LIGHT)


# ============================================================
# Slide 2: 目次
# ============================================================
slide = add_slide()
add_header(slide, "目次", "Agenda")
toc = [
    ("01", "プロジェクト概要", "AlphaScopeとは何か"),
    ("02", "開発の目的とゴール", "なぜ作るのか / 何を達成するか"),
    ("03", "スコープと方向性", "MVPで何をどこまで作るか"),
    ("04", "現状分析と着手方針", "どこから始めるか"),
    ("05", "クライアント3質問への回答", "予算・人件費・スケジュール"),
    ("06", "運用コスト見積", "サーバー・API等の月額"),
    ("07", "ビジネス優位性", "本プロジェクトの強み"),
    ("08", "マイルストーン", "段階別ゴールと作業内容"),
    ("09", "機能実装計画", "MVP / β版 / 本番の機能マトリクス"),
    ("10", "UI設計", "段階別の画面構成"),
]
y0 = Inches(1.3)
for i, (no, title, desc) in enumerate(toc):
    row_y = y0 + Inches(0.55) * i
    add_rect(slide, Inches(0.8), row_y, Inches(0.7), Inches(0.45), BLUE)
    add_text(slide, Inches(0.8), row_y, Inches(0.7), Inches(0.45),
             no, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(1.7), row_y, Inches(4), Inches(0.45),
             title, size=14, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(5.8), row_y, Inches(7), Inches(0.45),
             desc, size=12, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)
add_footer(slide, 2, 16)


# ============================================================
# Slide 3: プロジェクト概要
# ============================================================
slide = add_slide()
add_header(slide, "01. プロジェクト概要", "AlphaScope - AI投資情報分析SaaS")

# 左カラム: 概要文
add_rect(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.7), LIGHT)
add_text(slide, Inches(0.7), Inches(1.5), Inches(5.6), Inches(0.5),
         "■ サービス概要", size=16, bold=True, color=NAVY)
add_text(slide, Inches(0.7), Inches(2.05), Inches(5.6), Inches(2.0),
         "世界中のニュースから投資判断に必要な情報を\n"
         "AIが自動で収集・要約・スコアリングし、\n"
         "Discord等のチャネルに配信するSaaS。\n\n"
         "投資助言ではなく『意思決定のための情報提供』。",
         size=13, color=DARK)

add_text(slide, Inches(0.7), Inches(4.1), Inches(5.6), Inches(0.5),
         "■ ターゲット", size=16, bold=True, color=NAVY)
add_text(slide, Inches(0.7), Inches(4.65), Inches(5.6), Inches(1.0),
         "個人投資家 (初級〜中級層)\n"
         "→ 情報過多の中で『見るべき記事』を絞れない層",
         size=13, color=DARK)

add_text(slide, Inches(0.7), Inches(5.85), Inches(5.6), Inches(0.5),
         "■ 検証フェーズ", size=16, bold=True, color=NAVY)
add_text(slide, Inches(0.7), Inches(6.35), Inches(5.6), Inches(0.6),
         "テストユーザー10名にDiscord配信を行い、有用性を検証",
         size=13, color=DARK)

# 右カラム: フロー図
add_text(slide, Inches(7.0), Inches(1.4), Inches(5.8), Inches(0.5),
         "■ システムフロー", size=16, bold=True, color=NAVY)

flow_items = [
    ("ニュース収集", "RSS / API / スクレイピング\n3〜5媒体", BLUE),
    ("AI処理", "要約 / ジャンル分類 / スコアリング\n(OpenAI API)", ACCENT),
    ("配信制御", "ジャンル別頻度 / 重要度フィルタ\n緊急即時配信", GREEN),
    ("Discord配信", "ジャンル別CH / 全体CH / 緊急CH\n→ テストユーザー10名", NAVY),
]
fy = Inches(2.0)
fh = Inches(1.05)
gap = Inches(0.15)
for i, (label, desc, color) in enumerate(flow_items):
    y = fy + (fh + gap) * i
    add_rect(slide, Inches(7.0), y, Inches(5.8), fh, color)
    add_text(slide, Inches(7.2), y, Inches(2.0), fh, label,
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(9.2), y, Inches(3.5), fh, desc,
             size=11, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if i < len(flow_items) - 1:
        # 矢印 (▼)
        arrow_y = y + fh - Emu(20000)
        add_text(slide, Inches(9.5), arrow_y, Inches(0.8), gap,
                 "▼", size=12, bold=True, color=GRAY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_footer(slide, 3, 16)


# ============================================================
# Slide 4: 目的とゴール
# ============================================================
slide = add_slide()
add_header(slide, "02. 開発の目的とゴール", "なぜ作るのか / 何を達成するか")

# 3カラム (短期/中期/長期)
phases = [
    ("短期", "MVP完成と\n初期ユーザー検証",
     "・AIで投資情報を自動キュレーション\n"
     "・Discord配信を2週間で構築\n"
     "・10名のテストユーザーで実利用\n"
     "・有用性・配信頻度・通知UXを検証", BLUE),
    ("中期", "補助金申請用\nエビデンスの確立",
     "・MVPの利用実績・継続率を取得\n"
     "・事業性証明データを整備\n"
     "・申請通過 → 次フェーズの開発資金確保\n"
     "・β版開発の足場とする", ACCENT),
    ("長期", "投資判断支援SaaSとしての\n事業化",
     "・配信先をLINE/X等に拡張\n"
     "・有料プランの試行\n"
     "・全銘柄AI分析・自社AIモデル化\n"
     "・月額課金SaaSとして事業化", GREEN),
]
col_w = Inches(4.1)
gap = Inches(0.1)
start_x = Inches(0.5)
for i, (phase, title, body, color) in enumerate(phases):
    x = start_x + (col_w + gap) * i
    # ヘッダ
    add_rect(slide, x, Inches(1.3), col_w, Inches(0.6), color)
    add_text(slide, x, Inches(1.3), col_w, Inches(0.6),
             phase, size=16, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # タイトル
    add_rect(slide, x, Inches(1.9), col_w, Inches(1.1), LIGHT)
    add_text(slide, x, Inches(1.9), col_w, Inches(1.1),
             title, size=15, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 本文
    add_rect(slide, x, Inches(3.0), col_w, Inches(3.2), WHITE,
             line=RGBColor(0xBF, 0xBF, 0xBF))
    add_text(slide, x + Inches(0.15), Inches(3.15), col_w - Inches(0.3), Inches(2.9),
             body, size=12, color=DARK)

# MVP成功基準
add_rect(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.7), NAVY)
add_text(slide, Inches(0.7), Inches(6.4), Inches(12), Inches(0.7),
         "MVP成功基準: ①無人で1日2〜3回配信が安定 / ②10名テストユーザー稼働 / "
         "③緊急ニュース即時配信 / ④運用コスト月5万円以内",
         size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

add_footer(slide, 4, 16)


# ============================================================
# Slide 5: スコープと方向性
# ============================================================
slide = add_slide()
add_header(slide, "03. スコープと方向性", "MVPで何をどこまで作るか")

headers = ["カテゴリ", "対象", "MVPでの扱い", "理由・補足"]
rows = [
    ["データ収集", "経済ニュース 3〜5媒体", "○ 実装", "Yahoo!ファイナンス・東洋経済・ロイター等。RSS優先"],
    ["データ収集", "SNS (X / Reddit)", "− 対象外", "情報源を限定し品質優先。β版以降"],
    ["AI処理", "要約 + ジャンル分類", "○ 実装", "OpenAI API。プロンプト固定+JSON出力で品質安定"],
    ["AI処理", "重要度スコア (1〜5)", "○ 実装", "ジャンル別閾値で配信頻度を制御"],
    ["AI処理", "リード/ラグ分析", "△ 簡易のみ", "高度な相関分析はβ版。MVPは関連頻度UPのみ"],
    ["配信", "Discord Bot", "○ 実装", "ジャンル別CH + 全体CH + 緊急CH"],
    ["配信", "LINE / X / Webhook", "− 抽象化のみ", "Notifier基底クラスで本番差し替え可能に"],
    ["運用", "管理画面 (Web)", "− 対象外", "MVPはYAML+CLIで運用"],
    ["運用", "ユーザー管理", "− 対象外", "Discord参加=登録とみなす"],
    ["品質", "リトライ・ログ・JSON強制", "○ 実装", "外部API障害でもパイプラインを止めない"],
]
add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.7),
          headers, rows, col_ratios=[2, 4, 2.5, 6],
          row_height=Inches(0.5), body_font_size=11)
add_footer(slide, 5, 16)


# ============================================================
# Slide 6: 現状分析と着手方針
# ============================================================
slide = add_slide()
add_header(slide, "04. 現状分析と着手方針", "どこから始めるか / Week別作業")

# Week1, Week2, Week3 のタイムライン
weeks = [
    ("Week 1", "データ基盤 + AI処理",
     ["RSS/API取得モジュール",
      "スクレイピング (必要媒体のみ)",
      "正規化・重複排除",
      "PostgreSQL設計",
      "OpenAI APIラッパー (JSON強制)",
      "要約・分類・スコアプロンプト v1"], BLUE),
    ("Week 2", "Discord配信 + テスト",
     ["Discord Botセットアップ",
      "ジャンル別チャンネル設計",
      "定時配信スケジューラ",
      "緊急即時配信トリガ",
      "10名テストユーザー受入",
      "フィードバック収集"], ACCENT),
    ("Week 3 (任意)", "改善・最適化",
     ["プロンプト調整 (要約品質)",
      "配信頻度の調整",
      "コスト最適化 (キャッシュ等)",
      "障害復旧手順",
      "本番想定の負荷テスト",
      "運用ドキュメント整備"], GREEN),
]
col_w = Inches(4.1)
start_x = Inches(0.5)
for i, (week, theme, tasks, color) in enumerate(weeks):
    x = start_x + (col_w + Inches(0.1)) * i
    add_rect(slide, x, Inches(1.3), col_w, Inches(0.7), color)
    add_text(slide, x, Inches(1.3), col_w, Inches(0.7),
             week, size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, x, Inches(2.0), col_w, Inches(0.55), LIGHT)
    add_text(slide, x, Inches(2.0), col_w, Inches(0.55),
             theme, size=14, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, x, Inches(2.55), col_w, Inches(4.0), WHITE,
             line=RGBColor(0xBF, 0xBF, 0xBF))
    body = "\n".join(f"・{t}" for t in tasks)
    add_text(slide, x + Inches(0.2), Inches(2.7), col_w - Inches(0.4), Inches(3.7),
             body, size=12, color=DARK)

# 並行作業バー
add_rect(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5), NAVY)
add_text(slide, Inches(0.7), Inches(6.7), Inches(12), Inches(0.5),
         "並行: Discordチャンネル設計 (#要件整理 / #バグ報告 / #改善案) で混在を防止",
         size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

add_footer(slide, 6, 16)


# ============================================================
# Slide 7: クライアント3質問への回答 (1/2) - 予算・人件費
# ============================================================
slide = add_slide()
add_header(slide, "05. クライアント3質問への回答 (1/2)", "予算 / 人件費")

# Q1
add_rect(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(2.7), LIGHT)
add_rect(slide, Inches(0.5), Inches(1.3), Inches(2.5), Inches(2.7), BLUE)
add_text(slide, Inches(0.5), Inches(1.3), Inches(2.5), Inches(2.7),
         "Q1\n運用コスト\n概算",
         size=18, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(slide, Inches(3.2), Inches(1.45), Inches(9.4), Inches(0.5),
         "MVP段階 (テストユーザー10名)", size=14, bold=True, color=NAVY)
add_text(slide, Inches(3.2), Inches(1.85), Inches(9.4), Inches(0.5),
         "月額 15,000 〜 25,000円 が現実的レンジ",
         size=20, bold=True, color=ACCENT)
add_text(slide, Inches(3.2), Inches(2.45), Inches(9.4), Inches(1.5),
         "・サーバー(VPS/EC2): 月3,000〜8,000円\n"
         "・OpenAI API (gpt-4o-mini中心): 月5,000〜20,000円\n"
         "・ニュース取得: 無料tier+RSS+スクレイピング併用で月0円\n"
         "・Discord: 無料 / DB・ストレージ: 月数千円",
         size=12, color=DARK)

# Q2
add_rect(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.7), LIGHT)
add_rect(slide, Inches(0.5), Inches(4.1), Inches(2.5), Inches(2.7), ACCENT)
add_text(slide, Inches(0.5), Inches(4.1), Inches(2.5), Inches(2.7),
         "Q2\n人件費\n(開発費)",
         size=18, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(slide, Inches(3.2), Inches(4.25), Inches(9.4), Inches(0.5),
         "MVP一式 (1〜2週間 集中開発)", size=14, bold=True, color=NAVY)
add_text(slide, Inches(3.2), Inches(4.65), Inches(9.4), Inches(0.5),
         "20万円 (基本) / 30万円 (改善フェーズ込み)",
         size=20, bold=True, color=ACCENT)
add_text(slide, Inches(3.2), Inches(5.25), Inches(9.4), Inches(1.5),
         "・PM・要件整理・仕様調整・実装・テスト・初期ユーザー対応まで1名で完遂\n"
         "・募集記事ベースの基本価格は 20万円\n"
         "・要件擦り合わせ・改善ループまで含めるなら 30万円が現実的",
         size=12, color=DARK)

add_footer(slide, 7, 16)


# ============================================================
# Slide 8: クライアント3質問への回答 (2/2) - スケジュール
# ============================================================
slide = add_slide()
add_header(slide, "05. クライアント3質問への回答 (2/2)", "スケジュール")

add_rect(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.7), GREEN)
add_text(slide, Inches(0.7), Inches(1.3), Inches(12), Inches(0.7),
         "Q3: 実装までのスケジュール感 → 2週間で稼働、3週目で品質を整える",
         size=16, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# ガントチャート風
sched = [
    ("Phase 0: 要件整理", "0.5週", 0, 0.5, BLUE),
    ("Phase 1: データ収集基盤", "Week1 前半", 0.5, 1.0, BLUE),
    ("Phase 2: AI処理", "Week1 後半", 1.0, 1.0, BLUE),
    ("Phase 3: Discord配信", "Week2 前半", 2.0, 1.0, ACCENT),
    ("Phase 4: テストユーザー受入", "Week2 後半", 2.5, 0.5, ACCENT),
    ("Phase 5: 改善 (任意)", "Week3", 3.0, 1.0, GREEN),
]
chart_x = Inches(0.5)
chart_y = Inches(2.4)
label_w = Inches(3.5)
period_w = Inches(1.5)
bar_total_w = Inches(7.0)
unit_w = bar_total_w / 4  # 4週間想定 (0〜4)

# ヘッダ (Week目盛)
add_rect(slide, chart_x + label_w + period_w, chart_y, bar_total_w, Inches(0.4),
         NAVY)
for w in range(4):
    add_text(slide, chart_x + label_w + period_w + unit_w * w, chart_y,
             unit_w, Inches(0.4),
             f"Week {w+1}", size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

row_h = Inches(0.55)
for i, (name, period, start, dur, color) in enumerate(sched):
    y = chart_y + Inches(0.4) + row_h * i
    fill = WHITE if i % 2 == 0 else LIGHT
    add_rect(slide, chart_x, y, label_w + period_w + bar_total_w, row_h, fill,
             line=RGBColor(0xE0, 0xE0, 0xE0))
    add_text(slide, chart_x + Inches(0.1), y, label_w - Inches(0.1), row_h,
             name, size=11, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, chart_x + label_w, y, period_w, row_h,
             period, size=11, color=GRAY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bar_x = chart_x + label_w + period_w + Emu(int(unit_w * start))
    bar_w = Emu(int(unit_w * dur))
    add_rect(slide, bar_x + Emu(2000), y + Inches(0.08), bar_w - Emu(4000),
             row_h - Inches(0.16), color)

# 補足
add_rect(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7), LIGHT)
add_text(slide, Inches(0.7), Inches(6.5), Inches(12), Inches(0.7),
         "Phase 0は事前合意フェーズ。Week3は任意バッファ。"
         "Discordベースの非同期コミュニケーションで進行。",
         size=12, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

add_footer(slide, 8, 16)


# ============================================================
# Slide 9: 運用コスト見積 (3段階)
# ============================================================
slide = add_slide()
add_header(slide, "06. 運用コスト見積", "サーバー・API等の月額 ※開発費は含まず")

stages = [
    ("MVP", "テストユーザー10名\n1日2〜3回配信",
     "¥15,000 〜 ¥25,000",
     ["サーバー: VPS/EC2 small",
      "OpenAI API: gpt-4o-mini",
      "ニュース: 無料tier+RSS",
      "Discord: 無料",
      "DB: 同居 or RDS最小"], BLUE),
    ("β版", "ユーザー100名\n1日5〜10回配信\n10〜30媒体",
     "¥80,000 〜 ¥150,000",
     ["サーバー: t3.medium級",
      "OpenAI: 4o + 4o-mini",
      "ニュース: 商用API契約",
      "LINE / X 連携追加",
      "RDS + バックアップ"], ACCENT),
    ("本番", "ユーザー1,000名以上\n全銘柄分析\n複数チャネル",
     "¥200,000 〜 ¥500,000",
     ["ECS/Fargate or EKS",
      "自社モデル併用",
      "複数商用API",
      "Multi-AZ RDS",
      "監視: Datadog等"], NAVY),
]
col_w = Inches(4.1)
for i, (stage, desc, price, items, color) in enumerate(stages):
    x = Inches(0.5) + (col_w + Inches(0.1)) * i
    # ヘッダ
    add_rect(slide, x, Inches(1.3), col_w, Inches(0.7), color)
    add_text(slide, x, Inches(1.3), col_w, Inches(0.7),
             stage, size=20, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 規模
    add_rect(slide, x, Inches(2.0), col_w, Inches(1.0), LIGHT)
    add_text(slide, x, Inches(2.0), col_w, Inches(1.0),
             desc, size=12, color=DARK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 価格
    add_rect(slide, x, Inches(3.0), col_w, Inches(0.9), WHITE,
             line=color)
    add_text(slide, x, Inches(3.0), col_w, Inches(0.4),
             "月額", size=11, color=GRAY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x, Inches(3.35), col_w, Inches(0.55),
             price, size=20, bold=True, color=color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 内訳
    add_rect(slide, x, Inches(3.9), col_w, Inches(2.6), WHITE,
             line=RGBColor(0xBF, 0xBF, 0xBF))
    body = "\n".join(f"・{it}" for it in items)
    add_text(slide, x + Inches(0.2), Inches(4.05), col_w - Inches(0.4), Inches(2.4),
             body, size=12, color=DARK)

# コスト最適化
add_rect(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.6), NAVY)
add_text(slide, Inches(0.7), Inches(6.6), Inches(12), Inches(0.6),
         "コスト最適化: gpt-4o-mini主力 + 重要度=5のみ4oで再要約 → API費を1/3〜1/5に圧縮可",
         size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
add_footer(slide, 9, 16)


# ============================================================
# Slide 10: 運用コスト内訳 (MVP)
# ============================================================
slide = add_slide()
add_header(slide, "06. 運用コスト見積 - MVP内訳", "テストユーザー10名・1日2〜3回配信")

headers = ["項目", "推奨構成", "下限", "上限", "備考"]
rows = [
    ["サーバー", "AWS EC2 t3.small / VPS 2GB", "¥3,000", "¥8,000", "常時稼働"],
    ["DB", "PostgreSQL (RDS / VPS同居)", "¥0", "¥5,000", "MVPはVPS同居で削減"],
    ["ストレージ", "S3 / VPSローカル 10〜50GB", "¥0", "¥2,000", "圧縮で削減"],
    ["OpenAI API", "gpt-4o-mini中心 / 1日500〜2000記事", "¥5,000", "¥20,000", "キャッシュで圧縮可"],
    ["ニュースAPI", "NewsAPI無料tier + RSS + スクレイピング", "¥0", "¥0", "MVPは無料の範囲"],
    ["Discord", "Bot (無料)", "¥0", "¥0", "Botは無料"],
    ["監視", "Sentry無料tier等", "¥0", "¥2,000", "本格運用は本番で"],
    ["合計目安", "—", "¥8,000", "¥38,000", "実運用は¥15,000〜25,000に収まる"],
]
add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.0),
          headers, rows, col_ratios=[2, 5, 1.5, 1.5, 4],
          row_height=Inches(0.55), body_font_size=11)

add_rect(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.6), ACCENT)
add_text(slide, Inches(0.7), Inches(6.5), Inches(12), Inches(0.6),
         "→ 開発費(20万円)とは別に、月15,000〜25,000円の運用予算をご用意ください",
         size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
add_footer(slide, 10, 16)


# ============================================================
# Slide 11: ビジネス優位性
# ============================================================
slide = add_slide()
add_header(slide, "07. ビジネス優位性", "本プロジェクトの強み")

advantages = [
    ("速度", "OpenAI APIで2週間立ち上げ", "自社モデル開発不要 / 検証→改善ループが早い"),
    ("法務", "投資助言ではなく情報整理", "金商法登録不要のポジションで法的リスクを最小化"),
    ("拡張", "配信先抽象化", "Discord→LINE→X→アプリへバックエンド差し替えのみで対応"),
    ("市場", "個人投資家の情報過多を解く", "ジャンル別重要度フィルタの自動化は競合少"),
    ("資金", "補助金活用と段階的事業化", "MVP→補助金→β版→有料化の資金連動設計"),
    ("コスト", "従量×キャッシュ最適化", "ユーザー増にリニア比例しない構造を仕込む"),
]
col_w = Inches(4.0)
row_h = Inches(2.4)
gap = Inches(0.15)
for i, (key, title, desc) in enumerate(advantages):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + (col_w + gap) * col
    y = Inches(1.4) + (row_h + gap) * row
    # キーバッジ
    add_rect(slide, x, y, Inches(1.0), Inches(1.0), ACCENT)
    add_text(slide, x, y, Inches(1.0), Inches(1.0),
             key, size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # タイトル
    add_rect(slide, x + Inches(1.0), y, col_w - Inches(1.0), Inches(1.0), NAVY)
    add_text(slide, x + Inches(1.1), y, col_w - Inches(1.1), Inches(1.0),
             title, size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # 説明
    add_rect(slide, x, y + Inches(1.0), col_w, row_h - Inches(1.0), LIGHT)
    add_text(slide, x + Inches(0.15), y + Inches(1.05),
             col_w - Inches(0.3), row_h - Inches(1.1),
             desc, size=12, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

add_footer(slide, 11, 16)


# ============================================================
# Slide 12: マイルストーン
# ============================================================
slide = add_slide()
add_header(slide, "08. マイルストーン", "段階別ゴールと主要作業")

headers = ["フェーズ", "期間", "ゴール", "主な作業", "成果物"]
rows = [
    ["Phase 0\n要件整理", "0.5週", "媒体・ジャンル・配信ルール合意",
     "媒体リスト確定 / ジャンル設計 / 利用規約精査", "要件整理ドキュメント"],
    ["Phase 1\nデータ基盤", "Week1前半", "ニュースが自動でDB保存",
     "RSS/API取得 / スクレイピング / 重複排除", "収集パイプライン"],
    ["Phase 2\nAI処理", "Week1後半", "要約・分類・スコア出力",
     "OpenAI APIラッパー / プロンプトv1 / JSON強制", "AI処理モジュール"],
    ["Phase 3\nDiscord配信", "Week2前半", "ジャンル別/全体/緊急で配信",
     "Bot / スケジューラ / 配信フォーマット", "稼働中Discord Bot"],
    ["Phase 4\nテスト受入", "Week2後半", "10名がDiscordで配信受信",
     "サーバー設計 / 招待 / フィードバック収集", "稼働Discordサーバー"],
    ["Phase 5\n改善", "Week3 (任意)", "初期FBを反映し品質向上",
     "プロンプト調整 / 頻度調整 / コスト最適化", "改善版v1.1"],
    ["Phase 6\nβ版 (将来)", "1〜2ヶ月", "100名 / 配信先拡張 / 申請",
     "LINE/X追加 / 管理画面 / 有料tier試行", "β版 + 申請レポート"],
    ["Phase 7\n本番 (将来)", "3ヶ月以降", "全銘柄AI / 自社モデル化",
     "ファインチューニング / 課金 / SLA", "本番SaaS"],
]
add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.7),
          headers, rows, col_ratios=[2, 1.5, 3.5, 4, 3],
          row_height=Inches(0.7), body_font_size=10)
add_footer(slide, 12, 16)


# ============================================================
# Slide 13: 機能実装計画
# ============================================================
slide = add_slide()
add_header(slide, "09. 機能実装計画", "MVP / β版 / 本番の機能マトリクス")

headers = ["カテゴリ", "機能", "MVP", "β版", "本番", "備考"]
rows = [
    ["データ収集", "RSS / 公式API取得", "○", "○", "○", "Yahoo!ファイナンス・ロイター等"],
    ["データ収集", "スクレイピング", "△", "○", "○", "MVPは限定。利用規約遵守"],
    ["データ収集", "SNS取得 (X/Reddit)", "−", "△", "○", "情報源拡張"],
    ["AI処理", "要約 + ジャンル分類", "○", "○", "○", "gpt-4o-mini"],
    ["AI処理", "重要度スコア (1〜5)", "○", "○", "○", "ジャンル別閾値"],
    ["AI処理", "リスク抽出 / 関連分析", "△", "○", "○", "MVPは簡易"],
    ["AI処理", "自社モデル (FT)", "−", "−", "○", "本番で順次置換"],
    ["配信", "Discord Bot", "○", "○", "○", "ジャンル別/全体/緊急"],
    ["配信", "LINE / X / Webhook", "−", "○", "○", "Notifier抽象化済"],
    ["管理", "YAML/CLI 管理", "○", "○", "○", "MVPはこれで足りる"],
    ["管理", "Web管理画面 / ユーザー管理", "−", "○", "○", "β版以降"],
    ["管理", "課金システム", "−", "△", "○", "本番で本格化"],
    ["品質", "リトライ・JSON強制・ログ", "○", "○", "○", "基本対策"],
    ["品質", "監視 / アラート / オートスケール", "△", "○", "○", "段階的に強化"],
    ["分析", "配信反応・ユーザー行動分析", "−", "△", "○", "β版以降"],
]
add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5),
          headers, rows, col_ratios=[2, 4, 1, 1, 1, 4.5],
          row_height=Inches(0.36), body_font_size=10, header_font_size=11)

add_text(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
         "凡例: ○ = 実装する / △ = 部分的に実装 / − = 対象外",
         size=10, color=GRAY)
add_footer(slide, 13, 16)


# ============================================================
# Slide 14: UI設計 - MVP
# ============================================================
slide = add_slide()
add_header(slide, "10. UI設計 - MVP段階", "Discord中心 / 管理UIは最小")

# 左: Discord構成図
add_text(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(0.5),
         "■ Discord (テストユーザー10名向け)", size=15, bold=True, color=NAVY)

discord_channels = [
    ("# all", "全ジャンル横断の重要ニュース", BLUE),
    ("# 石油", "石油・エネルギー関連 (1h頻度)", BLUE),
    ("# 金融", "金融・為替関連", BLUE),
    ("# 地政学", "地政学・国際情勢", BLUE),
    ("# テック", "テック・半導体", BLUE),
    ("# urgent", "重要度=5の即時配信", ACCENT),
    ("# 要件整理", "クライアントとの仕様議論", GRAY),
    ("# バグ報告", "テストユーザーからのバグ報告", GRAY),
    ("# 改善案", "改善要望・フィードバック", GRAY),
]
ch_y = Inches(1.85)
for i, (name, desc, color) in enumerate(discord_channels):
    y = ch_y + Inches(0.42) * i
    add_rect(slide, Inches(0.5), y, Inches(1.5), Inches(0.38), color)
    add_text(slide, Inches(0.5), y, Inches(1.5), Inches(0.38),
             name, size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(2.0), y, Inches(4.5), Inches(0.38), LIGHT)
    add_text(slide, Inches(2.1), y, Inches(4.4), Inches(0.38),
             desc, size=10, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

# 右: 配信フォーマット例
add_text(slide, Inches(7.0), Inches(1.3), Inches(6.0), Inches(0.5),
         "■ 配信メッセージフォーマット (固定)", size=15, bold=True, color=NAVY)

add_rect(slide, Inches(7.0), Inches(1.85), Inches(6.0), Inches(3.4), WHITE,
         line=NAVY)
add_text(slide, Inches(7.2), Inches(2.0), Inches(5.6), Inches(0.4),
         "🔥 [重要度: ★★★★★]  #石油",
         size=14, bold=True, color=ACCENT)
add_text(slide, Inches(7.2), Inches(2.45), Inches(5.6), Inches(0.5),
         "中東情勢の緊迫化で原油先物が急騰",
         size=15, bold=True, color=DARK)
add_text(slide, Inches(7.2), Inches(3.0), Inches(5.6), Inches(1.5),
         "・WTI原油が前日比+3.2%で取引終了\n"
         "・市場はホルムズ海峡の供給リスクを警戒\n"
         "・関連: エネルギー株・防衛株が連動高",
         size=12, color=DARK)
add_text(slide, Inches(7.2), Inches(4.7), Inches(5.6), Inches(0.4),
         "🔗 元記事: https://example.com/...", size=11, color=BLUE)

# 管理方法
add_text(slide, Inches(7.0), Inches(5.4), Inches(6.0), Inches(0.5),
         "■ 管理 (Web UIなし)", size=15, bold=True, color=NAVY)
add_rect(slide, Inches(7.0), Inches(5.95), Inches(6.0), Inches(1.1), LIGHT)
add_text(slide, Inches(7.2), Inches(6.0), Inches(5.6), Inches(1.0),
         "・ジャンル別配信頻度・閾値: YAMLで管理\n"
         "・手動配信・再処理: CLI (python script)\n"
         "・状態確認: ログファイル / Sentry無料tier",
         size=11, color=DARK)

add_footer(slide, 14, 16)


# ============================================================
# Slide 15: UI設計 - β / 本番
# ============================================================
slide = add_slide()
add_header(slide, "10. UI設計 - β版 / 本番", "Web管理画面 + 複数チャネル + SaaS化")

# β版
add_rect(slide, Inches(0.5), Inches(1.3), Inches(6.1), Inches(0.6), ACCENT)
add_text(slide, Inches(0.5), Inches(1.3), Inches(6.1), Inches(0.6),
         "β版", size=18, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

beta_screens = [
    ("管理ダッシュボード", "配信状況・処理件数・コスト概算"),
    ("ジャンル管理", "ジャンル追加/編集/閾値設定"),
    ("プロンプト管理", "プロンプトのバージョン管理"),
    ("ユーザー管理", "ユーザー一覧・配信先設定"),
    ("ニュース閲覧 (任意)", "Discord履歴の補助 / 後追い閲覧"),
    ("個人設定", "ジャンル選択・配信先 (Discord/LINE/X)"),
]
for i, (name, desc) in enumerate(beta_screens):
    y = Inches(2.0) + Inches(0.65) * i
    add_rect(slide, Inches(0.5), y, Inches(2.5), Inches(0.55), LIGHT)
    add_text(slide, Inches(0.6), y, Inches(2.4), Inches(0.55),
             name, size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(3.0), y, Inches(3.6), Inches(0.55), WHITE,
             line=RGBColor(0xBF, 0xBF, 0xBF))
    add_text(slide, Inches(3.1), y, Inches(3.5), Inches(0.55),
             desc, size=11, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

# 本番
add_rect(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(0.6), NAVY)
add_text(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(0.6),
         "本番", size=18, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

prod_screens = [
    ("LP / 申込ページ", "サービス紹介・料金プラン・登録導線"),
    ("ユーザー Dashboard", "重要ニュース・銘柄WL・履歴"),
    ("分析レポート (有料)", "週次/月次のジャンル別動向"),
    ("個人設定", "ジャンル・配信先・通知時間・プラン"),
    ("運営: 課金管理", "Stripe等との連携"),
    ("運営: AI出力評価", "プロンプト品質モニタリング"),
]
for i, (name, desc) in enumerate(prod_screens):
    y = Inches(2.0) + Inches(0.65) * i
    add_rect(slide, Inches(6.8), y, Inches(2.5), Inches(0.55), LIGHT)
    add_text(slide, Inches(6.9), y, Inches(2.4), Inches(0.55),
             name, size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(9.3), y, Inches(3.5), Inches(0.55), WHITE,
             line=RGBColor(0xBF, 0xBF, 0xBF))
    add_text(slide, Inches(9.4), y, Inches(3.4), Inches(0.55),
             desc, size=11, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

# 設計方針
add_rect(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7), NAVY)
add_text(slide, Inches(0.7), Inches(6.5), Inches(12), Inches(0.7),
         "設計方針: MVPは画面ゼロでDiscordに集約 → β版で管理画面 → 本番でユーザー向けWeb / 配信先抽象化はMVPから仕込む",
         size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
add_footer(slide, 15, 16)


# ============================================================
# Slide 16: クロージング
# ============================================================
slide = add_slide()
add_rect(slide, 0, 0, SW, SH, NAVY)
add_rect(slide, 0, Inches(2.4), SW, Inches(0.08), ACCENT)

add_text(slide, Inches(0.8), Inches(1.0), SW - Inches(1.6), Inches(0.6),
         "次のステップ", size=18, color=LIGHT)
add_text(slide, Inches(0.8), Inches(1.5), SW - Inches(1.6), Inches(0.8),
         "Let's build AlphaScope.",
         size=44, bold=True, color=WHITE)

steps = [
    ("01", "対象媒体 (3〜5サイト) のご選定 / ご共有"),
    ("02", "ジャンル分類・配信頻度ルールの初期合意"),
    ("03", "Lancers経由でご契約 → Discord開設 → 着手"),
    ("04", "Week2末に10名テストユーザー受入準備完了"),
]
sy = Inches(2.9)
for i, (no, txt) in enumerate(steps):
    y = sy + Inches(0.7) * i
    add_rect(slide, Inches(0.8), y, Inches(0.7), Inches(0.5), ACCENT)
    add_text(slide, Inches(0.8), y, Inches(0.7), Inches(0.5),
             no, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(1.7), y, SW - Inches(2.5), Inches(0.5),
             txt, size=14, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

add_text(slide, Inches(0.8), Inches(6.6), SW - Inches(1.6), Inches(0.5),
         "ご質問・ご要望はDiscordにてお気軽にお寄せください。",
         size=14, color=LIGHT)

# ============================================================
# 保存
# ============================================================
prs.save(OUT_PATH)
